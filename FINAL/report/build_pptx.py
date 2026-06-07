"""Build the GRACE Track-1 presentation (PPTX) from report_data.json + figures/.

Dark 16:9 deck (~28 slides) matching the figure theme so charts blend in.
Run from FINAL/:  python report/build_pptx.py
"""
from __future__ import annotations

import json
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
DATA = json.load(open(os.path.join(HERE, "report_data.json"), encoding="utf-8"))
LB = DATA["leaderboard"]
ENS = DATA["ensemble"]
DS = DATA["dataset"]
META = DATA["run_meta"]
ALL9 = ENS["dev_scores"]["all9"]["official_score"]
TOP3 = ENS["dev_scores"]["top3"]["official_score"]
GREEDY = ENS["dev_scores"]["greedy"]["official_score"]
BEST_SEED = max(max(r["seed_macro"].values()) for r in LB)
GAIN = (ALL9 - LB[0]["mean_macro"]) * 100

BG = RGBColor.from_string("0D1117")
PANEL = RGBColor.from_string("161B22")
FG = RGBColor.from_string("E6EDF3")
MUTED = RGBColor.from_string("8B949E")
ACCENT = RGBColor.from_string("7EE787")
ACCENT2 = RGBColor.from_string("58A6FF")
ACCENT3 = RGBColor.from_string("D2A8FF")
GOLD = RGBColor.from_string("E3B341")
WARN = RGBColor.from_string("F0883E")
FONT = "Segoe UI"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

IDX = 0
TOTAL = 28


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size, color=FG, bold=False, align=PP_ALIGN.LEFT, first=False,
         space_after=6, italic=False, font=FONT, level=0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return p


def multi(tf, segments, align=PP_ALIGN.LEFT, first=False, space_after=6):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    for text, opts in segments:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", 13))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.name = FONT
        r.font.color.rgb = opts.get("color", FG)
    return p


def bar(s, x, y, w, h, color):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def rrect(s, x, y, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def footer(s):
    tf = box(s, 0.5, 7.06, 9, 0.3)
    para(tf, "GRACE @ IberLEF 2026  ·  Track 1: Evidence Component Detection (ES)", 9, MUTED, first=True)
    tf2 = box(s, SW - 1.6, 7.06, 1.1, 0.3)
    para(tf2, f"{IDX} / {TOTAL}", 9, MUTED, align=PP_ALIGN.RIGHT, first=True)


def header(s, title, subtitle=None):
    bar(s, 0.5, 0.55, 0.14, 0.62, ACCENT)
    tf = box(s, 0.78, 0.42, 11.9, 1.05)
    para(tf, title, 29, FG, bold=True, first=True, space_after=2)
    if subtitle:
        para(tf, subtitle, 14, MUTED)


def begin(title=None, subtitle=None):
    global IDX
    IDX += 1
    s = slide()
    if title:
        header(s, title, subtitle)
    return s


def add_fig(s, name, x, y, max_w, max_h, center_x=True):
    path = os.path.join(FIG, name)
    w, h = Image.open(path).size
    ratio = min((max_w * 96) / w, (max_h * 96) / h)
    fw, fh = w * ratio / 96.0, h * ratio / 96.0
    if center_x:
        x = x + (max_w - fw) / 2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(fw), Inches(fh))


def fig_slide(title, subtitle, fig_name, note=None):
    s = begin(title, subtitle)
    add_fig(s, fig_name, 0.6, 1.5, 12.1, 5.05)
    if note:
        tf = box(s, 0.78, 6.62, 11.8, 0.4)
        para(tf, note, 12, ACCENT, italic=True, first=True)
    footer(s)
    return s


def cards(title, subtitle, items, top=1.62, ch=None, gap=0.16, body_size=13, bottom=6.95):
    s = begin(title, subtitle)
    n = len(items)
    ch = ch or (bottom - top - gap * (n - 1)) / n  # fit all cards above the footer
    y = top
    for ttl, body, col in items:
        rrect(s, 0.6, y, 12.1, ch, PANEL)
        bar(s, 0.6, y, 0.1, ch, col)
        tf = box(s, 0.95, y + 0.12, 11.55, ch - 0.2)
        para(tf, ttl, 15.5, col, bold=True, first=True, space_after=3)
        para(tf, body, body_size, FG)
        y += ch + gap
    footer(s)
    return s


def two_col(title, subtitle, left_head, left_items, right_head, right_items):
    s = begin(title, subtitle)
    for head, items, x in ((left_head, left_items, 0.6), (right_head, right_items, 6.75)):
        tf = box(s, x + 0.05, 1.55, 5.9, 0.4)
        para(tf, head, 15, ACCENT, bold=True, first=True)
        y = 2.12
        for ttl, body in items:
            rrect(s, x, y, 5.95, 1.07, PANEL)
            bar(s, x, y, 0.08, 1.07, ACCENT2)
            t = box(s, x + 0.26, y + 0.09, 5.55, 0.9)
            para(t, ttl, 13, ACCENT2, bold=True, first=True, space_after=2)
            para(t, body, 11.5, FG)
            y += 1.2
    footer(s)
    return s


def table_slide(title, subtitle, headers, rows, col_w, highlight=None, note=None, fontsize=11):
    s = begin(title, subtitle)
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = sum(col_w)
    x = (SW - total_w) / 2
    top = 1.7
    h = min(4.6, 0.42 * nrows)
    gt = s.shapes.add_table(nrows, ncols, Inches(x), Inches(top), Inches(total_w), Inches(h)).table
    gt.first_row = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for j, htext in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string("1A7F37")
        c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext; r.font.size = Pt(fontsize); r.font.bold = True
        r.font.color.rgb = FG; r.font.name = FONT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if (i % 2) else RGBColor.from_string("1B222B")
            if highlight is not None and i - 1 == highlight:
                c.fill.fore_color.rgb = RGBColor.from_string("23311f")
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val); r.font.size = Pt(fontsize)
            r.font.name = FONT
            r.font.bold = (highlight is not None and i - 1 == highlight and j == 0)
            r.font.color.rgb = GOLD if (highlight is not None and i - 1 == highlight) else FG
    if note:
        tf = box(s, 0.78, top + h + 0.2, 11.8, 0.6)
        para(tf, note, 12, ACCENT, italic=True, first=True)
    footer(s)
    return s


# ============================================================ 1 TITLE
global_ = begin()
s = global_
bar(s, 0, 0, SW, 0.16, ACCENT)
bar(s, 0, SH - 0.16, SW, 0.16, ACCENT)
tf = box(s, 1.0, 2.0, 11.3, 2.6)
para(tf, "Argument Mining in Spanish Clinical Trials", 40, FG, bold=True, first=True, space_after=8)
para(tf, "A 27-Run Transformer Benchmark + Ensemble for Evidence Component Detection",
     20, ACCENT, space_after=18)
para(tf, "GRACE @ IberLEF 2026   ·   Track 1   ·   Subtask 1", 16, MUTED)
tf2 = box(s, 1.0, 5.55, 11.3, 1.0)
multi(tf2, [("Best single: ", {"color": MUTED}), (f"mDeBERTa-v3  {LB[0]['mean_macro']:.3f}", {"bold": True, "color": FG}),
            ("      All-9 ensemble: ", {"color": MUTED}), (f"{ALL9:.3f}", {"bold": True, "color": GOLD}),
            ("      Blind: ", {"color": MUTED}), (f"{ENS['blind']['all9']['entities']:,} components", {"bold": True, "color": FG})],
      first=True)
footer(s)

# ============================================================ 2 AGENDA
s = begin("Agenda", "From a small, imbalanced corpus to a validated submission")
agenda = [
    ("01", "Problem & Data", "The task, an annotated example, the corpus and its class imbalance", ACCENT2),
    ("02", "Method", "Sentence-level 4-class framing, two-stage runner, the 9 backbones", ACCENT),
    ("03", "Infrastructure", "The GPU journey (Kaggle → Azure → Modal) and the parallel sweep", ACCENT3),
    ("04", "Results", "Leaderboard, per-class behaviour, seed stability, training dynamics", GOLD),
    ("05", "Ensemble & Submission", "Probability averaging, honest selection, the blind-test output", ACCENT),
    ("06", "Lessons & Outlook", "What we explored, engineering lessons, limitations, next steps", ACCENT2),
]
y = 1.7
for num, ttl, body, col in agenda:
    rrect(s, 0.6, y, 12.1, 0.78, PANEL)
    bar(s, 0.6, y, 0.1, 0.78, col)
    tnum = box(s, 0.85, y + 0.13, 1.0, 0.55)
    para(tnum, num, 22, col, bold=True, first=True)
    tf = box(s, 1.95, y + 0.07, 10.6, 0.7)
    para(tf, ttl, 15, col, bold=True, first=True, space_after=1)
    para(tf, body, 12, MUTED)
    y += 0.9
footer(s)

# ============================================================ 3 MOTIVATION
cards("Why Argument Mining in Clinical Trials?",
      "Turning prose evidence into structured, machine-readable argumentation",
      [("Evidence-based medicine runs on RCTs", "Randomised controlled trials are the gold standard of clinical "
        "evidence, but their conclusions are locked in free-text abstracts — hard to search, compare, or synthesise at scale.", ACCENT2),
       ("Arguments have structure", "Each abstract builds an argument: Premises (results) support Claims "
        "(interpretations) that roll up into a MajorClaim (the take-home conclusion). Recovering that structure powers "
        "automated evidence synthesis and systematic reviews.", ACCENT),
       ("Spanish clinical NLP is under-resourced", "Most argument-mining work targets English. GRACE provides a "
        "Spanish RCT benchmark — exactly where domain-specific models and careful methodology can make a difference.", ACCENT3),
       ("Our goal", "A rigorous, reproducible benchmark of Spanish/clinical encoders for evidence component "
        "detection, plus an ensemble strong enough to submit.", GOLD)])

# ============================================================ 4 TASK
s = begin("The Task — Evidence Component Detection",
          "Classify each segment of a Spanish RCT abstract by argumentative role")
labels = [("Premise", "Reported evidence,\nstatistical results", ACCENT2),
          ("Claim", "Interpretation of\nthe findings", ACCENT),
          ("MajorClaim", "The overall\nconclusion", GOLD),
          ("None", "Background /\nmethods text", MUTED)]
x = 0.6
for name, desc, col in labels:
    rrect(s, x, 1.7, 2.92, 1.6, PANEL)
    bar(s, x, 1.7, 2.92, 0.1, col)
    tf = box(s, x + 0.2, 2.0, 2.55, 1.2)
    para(tf, name, 17, col, bold=True, first=True, space_after=4)
    para(tf, desc, 12, MUTED)
    x += 3.04
tf = box(s, 0.6, 3.7, 12.1, 2.6)
para(tf, "Our framing: sentence-level 4-class classification.", 16, FG, bold=True, first=True, space_after=8)
para(tf, "• Segment each abstract into sentences with spaCy (es_core_news_sm).", 13.5, FG, space_after=4)
para(tf, "• Classify each sentence in the context of its previous sentence:  prev_sentence [SEP] sentence  (≤192 tokens).", 13.5, FG, space_after=4)
para(tf, "• A 0.50 probability override promotes a segment to the rare MajorClaim class.", 13.5, FG, space_after=4)
para(tf, "• Official ranking metric: strict macro-F1 — predicted spans must match gold character offsets exactly.", 13.5, FG, space_after=4)
footer(s)

# ============================================================ 5 ANNOTATED EXAMPLE
fig_slide("What the Model Labels", "Real dev examples of each argumentative role",
          "09_annotated_example.png",
          note="Premise → Claim → MajorClaim form the argumentative chain; everything else is None.")

# ============================================================ 6 DATASET
fig_slide("Dataset & the Core Challenge", "Small, severely imbalanced — MajorClaim is the bottleneck",
          "06_data_overview.png",
          note=f"350 train / 50 dev / {DS['docs']['blind']:,} blind abstracts · MajorClaim = 2.8% of train, only {DS['entities']['dev']['MajorClaim']} dev spans.")

# ============================================================ 7 EVAL METHODOLOGY
fig_slide("Evaluation Methodology", "Strict vs relaxed span matching, and macro-F1 aggregation",
          "11_eval_methodology.png",
          note="We optimise and report the official strict macro-F1 throughout.")

# ============================================================ 8 ARCHITECTURE
fig_slide("System Architecture", "One pipeline, two stages, 27 runs in parallel",
          "10_pipeline.png",
          note="Stage 1 selects the best epoch on held-out dev; Stage 2 refits on all data for the submission.")

# ============================================================ 9 INPUT REPRESENTATION
s = begin("Input Representation", "Context-aware sentence classification")
rrect(s, 0.6, 1.6, 12.1, 1.72, PANEL)
tf = box(s, 0.95, 1.72, 11.5, 1.5)
para(tf, "Each sentence is classified together with the one before it:", 14, FG, bold=True, first=True, space_after=10)
multi(tf, [("[CLS] ", {"color": MUTED, "size": 13}),
           ("La calidad de vida fue similar en los dos grupos.", {"color": ACCENT2, "size": 13, "bold": True}),
           (" [SEP]", {"color": GOLD, "size": 13, "bold": True}),
           ("    ← previous sentence (context)", {"color": ACCENT2, "size": 10.5, "italic": True})], space_after=6)
multi(tf, [("El resultado de estos pacientes sigue siendo insatisfactorio.", {"color": ACCENT, "size": 13, "bold": True}),
           (" [SEP]", {"color": GOLD, "size": 13, "bold": True}),
           ("    ← target sentence (the one we label)", {"color": ACCENT, "size": 10.5, "italic": True})], space_after=0)
items = [("Why previous-sentence context?", "Argument role is relational — a Claim often follows the Premises that "
         "support it. The preceding sentence gives the encoder that discourse cue.", ACCENT2),
         ("Segmentation", "spaCy es_core_news_sm splits abstracts into sentences; segments shorter than 10 characters "
          "are dropped as noise.", ACCENT),
         ("Length budget", "max_len 192 tokens comfortably covers a sentence pair in clinical Spanish without "
          "truncation in the vast majority of cases.", ACCENT3)]
y = 3.5
for ttl, body, col in items:
    rrect(s, 0.6, y, 12.1, 1.0, PANEL)
    bar(s, 0.6, y, 0.1, 1.0, col)
    t = box(s, 0.95, y + 0.1, 11.5, 0.85)
    para(t, ttl, 14, col, bold=True, first=True, space_after=2)
    para(t, body, 12, FG)
    y += 1.12
footer(s)

# ============================================================ 10 MODEL ROSTER
DOMAIN = {
    "mdeberta_v3": ("General web", "Multilingual", "base"),
    "bsc_bio_ehr": ("Biomedical + clinical EHR", "Spanish", "base"),
    "mrbert_es": ("General (ModernBERT)", "Spanish + English", "base"),
    "rigoberta_clinical": ("Clinical", "Spanish", "large"),
    "mrbert_biomed": ("Biomedical", "Spanish + English", "base"),
    "roberta_clinical": ("Biomedical + clinical", "Spanish", "base"),
    "xlmr_large": ("General web (100 langs)", "Multilingual", "large"),
    "beto": ("General web", "Spanish", "base"),
    "beto_galen": ("Clinical (Galén)", "Spanish", "base"),
}
rows = [[r["pretty"], *DOMAIN[r["prefix"]], f"{r['mean_macro']:.3f}"] for r in LB]
table_slide("The 9 Backbones", "General, biomedical, and clinical encoders — Spanish and multilingual",
            ["Model", "Pretraining domain", "Language", "Size", "Dev macro-F1"],
            rows, col_w=[3.2, 3.4, 2.2, 1.3, 2.0], highlight=0,
            note="Selected to span the pretraining spectrum — does clinical pretraining beat a strong multilingual general model? (It doesn't, quite.)",
            fontsize=11)

# ============================================================ 11 TRAINING SETUP
two_col("Training Setup", "Identical recipe across all 9 backbones for a fair comparison",
        "Handling class imbalance",
        [("Weighted cross-entropy", "Per-class weights with MajorClaim set to 10× to counter its 2.8% frequency."),
         ("Label smoothing 0.15", "Reduces over-confidence on the dominant Premise class and improves calibration."),
         ("0.50 MajorClaim override", "Promote a segment to MajorClaim when its probability clears 0.50 — recovers recall on the rare class.")],
        "Optimisation",
        [("Differential LR", "Body 1e-5, classifier head 1e-4 (AdamW, weight decay 0.05)."),
         ("Cosine schedule", "100 warmup steps, batch size 16, up to 15 epochs."),
         ("Early stopping", "Patience 4 on dev strict macro-F1; best epoch restored (load-best-at-end).")])

# ============================================================ 12 ENGINEERING JOURNEY
fig_slide("Infrastructure — The Path to a Working Sweep",
          "Three platforms, two dead-ends, one win on a tight GPU budget",
          "12_engineering_journey.png",
          note="Live monitoring was essential — it exposed the silent P100 hang that looked like normal training.")

# ============================================================ 13 MODAL PANEL
fig_slide("Infrastructure — One Parallel Sweep",
          "The full 9×3 grid trained in a single Modal launch",
          "07_run_panel.png")

# ============================================================ 14 BF16 FIX
s = begin("Deep-Dive — The bf16 Fix", "How a precision setting turned the worst model into the best")
steps = [("Symptom", "Under fp32, mDeBERTa-v3 produced grad_norm = NaN within the first epoch and collapsed to "
          "0.0 macro-F1 — effectively unusable.", WARN),
         ("Diagnosis", "mDeBERTa-v3's disentangled attention is numerically sensitive; fp32 mixed with the training "
          "dynamics overflowed. A common Kaggle workaround was pinning an older transformers version.", ACCENT2),
         ("Fix", "Train in bf16 on Ampere (A10G). bf16's wider dynamic range stabilises the gradients without loss "
          "scaling — no version pinning needed.", ACCENT),
         ("Result", f"mDeBERTa-v3 became the #1 single backbone at {LB[0]['mean_macro']:.3f} ± {LB[0]['std']:.3f} "
          f"dev macro-F1. The right precision mattered more than the architecture.", GOLD)]
y = 1.65
for ttl, body, col in steps:
    rrect(s, 0.6, y, 12.1, 1.16, PANEL)
    bar(s, 0.6, y, 0.1, 1.16, col)
    t = box(s, 0.95, y + 0.11, 11.5, 1.0)
    para(t, ttl, 15, col, bold=True, first=True, space_after=3)
    para(t, body, 12.5, FG)
    y += 1.28
footer(s)

# ============================================================ 15-18 RESULTS FIGURES
fig_slide("Results — Model Leaderboard", "Dev strict macro-F1, mean ± std over 3 seeds",
          "01_leaderboard.png",
          note=f"Clinical pretraining helps, but multilingual mDeBERTa-v3 wins; the ensemble line ({ALL9:.3f}) clears every single model.")
fig_slide("Results — Per-Class Breakdown", "Premise is easy; MajorClaim is the universal weak spot",
          "02_per_class_f1.png",
          note="0.82–0.90 on Premise vs 0.24–0.43 on MajorClaim — a data-scarcity ceiling, not a model defect.")
fig_slide("Results — Seed Stability", "Why we report mean ± std, not a single lucky run",
          "03_seed_variance.png",
          note="Seed spread reaches ~0.07 macro-F1 — single-seed leaderboards are unreliable on data this small.")
fig_slide("Results — Training Dynamics", "Fast convergence; best epoch kept via load-best-at-end",
          "05_training_curves.png",
          note="The five strongest models peak within 4–7 epochs (seed 42); early stopping (patience 4) prevents overfitting on 350 docs.")

# ============================================================ 19 ERROR ANALYSIS
cards("Error Analysis — Where Models Struggle",
      "Three recurring failure modes, all concentrated in the harder classes",
      [("MajorClaim ↔ Claim confusion", "The two are linguistically similar (both interpretive); with only 9 dev / "
        "64 train MajorClaim spans, models under-recall it and frequently label it Claim instead.", GOLD),
       ("Strict boundaries penalise near-misses", "Sentence-level spans sometimes differ from gold clause-level "
        "annotations by a few tokens — correct in spirit, but counted wrong under exact-offset strict matching.", WARN),
       ("Premise is near-saturated", "At 0.82–0.90 F1 across all models, Premise offers little headroom — gains must "
        "come from the rare classes, where data is scarcest.", ACCENT2)],
      body_size=13)

# ============================================================ 20 ENSEMBLE METHOD
fig_slide("Ensemble — Methodology", "Probability averaging over all 27 runs",
          "13_ensemble_method.png",
          note="Average softmax outputs (mean over seeds, then over models), then argmax with the 0.50 MajorClaim override.")

# ============================================================ 21 ENSEMBLE RESULTS
s = fig_slide("Ensemble — Honest Model Selection",
              "All-9 averaging beats every single model AND dev-selected top-3",
              "04_ensemble_gain.png",
              note=f"All-9 ({ALL9:.4f}) > top-3 ({TOP3:.4f}): the zero-peek choice is also the best → submit all-9, no selection bias.")

# ============================================================ 22 BLIND GENERATION
s = begin("Blind-Test Submission — Generation & Validation",
          "Applying the all-9 ensemble to 2,474 unseen abstracts")
b = ENS["blind"]["all9"]
kpis = [(f"{DS['docs']['blind']:,}", "abstracts", ACCENT2), (f"{DS['blind_sentences']:,}", "sentences", ACCENT3),
        (f"{b['entities']:,}", "components", ACCENT), ("0", "offset mismatches", GOLD)]
x = 0.6
for val, lab, col in kpis:
    rrect(s, x, 1.65, 2.92, 1.25, PANEL)
    bar(s, x, 1.65, 2.92, 0.09, col)
    tf = box(s, x + 0.1, 1.85, 2.72, 1.0)
    para(tf, val, 26, FG, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
    para(tf, lab, 12, MUTED, align=PP_ALIGN.CENTER)
    x += 3.04
checks = [("Stage-2 logits", "Predictions come from models refit on train+dev — the full-data fit, not the dev-selection fit.", ACCENT),
          ("Offset validation", "Every predicted span satisfies raw_text[start:end] == text after trimming spaCy's trailing whitespace (135 → 0).", ACCENT2),
          ("Sanity check", "Premise/Claim proportions match the training prior; MajorClaim is predicted conservatively (1.4% vs 2.8%) — no degenerate single-class collapse.", ACCENT3)]
y = 3.2
for ttl, body, col in checks:
    rrect(s, 0.6, y, 12.1, 1.02, PANEL)
    bar(s, 0.6, y, 0.1, 1.02, col)
    t = box(s, 0.95, y + 0.1, 11.5, 0.85)
    para(t, ttl, 14, col, bold=True, first=True, space_after=2)
    para(t, body, 12, FG)
    y += 1.14
footer(s)

# ============================================================ 23 BLIND DISTRIBUTION
fig_slide("Blind-Test Submission — Predicted Distribution",
          "All-9 ensemble output across the unseen abstracts",
          "08_blind_distribution.png",
          note="Premise/Claim mirror the training prior; MajorClaim is predicted more conservatively — the ensemble generalises without collapsing to one class.")

# ============================================================ 24 WHAT WE EXPLORED
cards("Beyond the Core Benchmark — What We Explored",
      "The sentence-level recipe was selected after a broader search",
      [("Architecture variants", "Focal loss, a CRF decoding layer, and a BiLSTM head were trialled on the encoder "
        "baseline before settling on the simpler, more reliable sentence-level classifier.", ACCENT2),
       ("Relation classification (Track 2)", "An NLI-style relation classifier for Support / Attack / Partial-Attack "
        "between components — the natural next step on top of detected evidence.", ACCENT),
       ("Data strategies", "Rare-class oversampling, external clinical corpora, and LLM-based synthetic augmentation "
        "to attack the MajorClaim scarcity directly.", ACCENT3),
       ("Why the sentence-level recipe won", "It was the most robust and reproducible across backbones and seeds — "
        "the right foundation for a fair 27-run benchmark and a trustworthy ensemble.", GOLD)],
      body_size=12.5)

# ============================================================ 25 LESSONS
cards("Engineering Lessons",
      "Four takeaways that generalise beyond this task",
      [("Precision can unlock a model", "Ampere bf16 turned mDeBERTa-v3 from a NaN collapse into the #1 backbone — "
        "the numerical setup mattered more than the architecture choice.", ACCENT),
       ("Ensembling is the cheap, honest win", f"Averaging 27 distributions adds +{GAIN:.1f} pts over the best single "
        f"model with zero extra training, and — chosen as all-9 — carries no selection bias.", GOLD),
       ("Report distributions, not lucky seeds", "Up to ~0.07 seed spread on 350 docs makes single-seed rankings "
        "unreliable; mean ± std over 3 seeds is the defensible comparison.", ACCENT2),
       ("Monitor training live", "The silent P100 hang looked like normal training — only live log monitoring "
        "revealed it. Trust, but verify your GPU is actually working.", ACCENT3)],
      body_size=12.5)

# ============================================================ 26 LIMITATIONS & FUTURE
two_col("Limitations & Future Work", "Honest scope and the road ahead",
        "Limitations",
        [("Sentence-level granularity", "Cannot recover sub-sentence clause spans that strict matching rewards."),
         ("Tiny dev set", "50 docs / 9 MajorClaim spans → wide confidence intervals on the rare class."),
         ("Dev-only estimate", "Blind labels are unreleased; 0.743 is our dev estimate, not a leaderboard score.")],
        "Future work",
        [("Clause-level segmentation", "Predict intrasentence spans to align with gold boundaries."),
         ("MajorClaim augmentation", "Targeted synthetic / contrastive data for the weakest class."),
         ("Track 2 relations", "Support / Attack / Partial-Attack detection on the detected components.")])

# ============================================================ 27 TAKEAWAYS
s = begin("Key Takeaways", None)
big = [("0.743", "dev strict macro-F1 from the bias-free all-9 ensemble", GOLD),
       (f"+{GAIN:.1f} pts", "ensemble gain over the best single model — no extra training", ACCENT),
       ("27 / 27", "runs converged in one parallel Modal launch", ACCENT2)]
x = 0.6
for val, lab, col in big:
    rrect(s, x, 1.6, 3.93, 1.55, PANEL)
    bar(s, x, 1.6, 3.93, 0.1, col)
    tf = box(s, x + 0.2, 1.85, 3.55, 1.2)
    para(tf, val, 34, col, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
    para(tf, lab, 12, MUTED, align=PP_ALIGN.CENTER)
    x += 4.07
tf = box(s, 0.6, 3.55, 12.1, 3.0)
para(tf, "In one sentence:", 15, FG, bold=True, first=True, space_after=8)
para(tf, "A disciplined 27-run benchmark plus a zero-peek ensemble delivers a strong, reproducible Spanish "
         "clinical argument-mining system — and the remaining headroom is a data problem (MajorClaim), not a "
         "modelling one.", 15, FG, space_after=14)
para(tf, "Methodology that travels:", 14, ACCENT, bold=True, space_after=6)
para(tf, "• Get the numerics right before blaming the architecture.    • Prefer bias-free ensembling to "
         "dev-tuned selection.    • Rank on mean ± std, never a lucky seed.", 13, FG)
footer(s)

# ============================================================ 28 THANK YOU
s = begin()
bar(s, 0, 0, SW, 0.16, ACCENT)
bar(s, 0, SH - 0.16, SW, 0.16, ACCENT)
tf = box(s, 1.0, 2.5, 11.3, 2.0)
para(tf, "Thank you", 44, FG, bold=True, first=True, space_after=10)
para(tf, "Questions & discussion welcome", 20, ACCENT, space_after=18)
tf2 = box(s, 1.0, 5.2, 11.3, 1.2)
para(tf2, "GRACE @ IberLEF 2026 · Track 1 — Evidence Component Detection (ES)", 14, MUTED, first=True, space_after=4)
para(tf2, "Reproducible pipeline: modal_sweep.py → modal_ensemble.py → report/ (figures · DOCX · PPTX)", 12, MUTED)
footer(s)

out = os.path.join(HERE, "GRACE_Track1_IberLEF2026.pptx")
prs.save(out)
n = len(prs.slides._sldIdLst)
print(f"saved {out}  ({n} slides)")
assert n == TOTAL, f"slide count {n} != TOTAL {TOTAL}"
